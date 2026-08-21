from __future__ import annotations

"""
Azure-integrated OCR engine — serverless OCR & layout parsing with digital PDF cost optimization.
"""

import csv
import io
import json
import logging
import os
import re
import tempfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

from .config import settings

logger = logging.getLogger("ocr.engine")

PageResult = tuple[int, str, float | None]

_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp", ".gif", ".heic", ".heif"}
_PDF_EXTENSIONS = {".pdf"}
_DOCX_EXTENSIONS = {".docx", ".doc"}
_EXCEL_EXTENSIONS = {".xlsx", ".xls"}
_PPTX_EXTENSIONS = {".pptx", ".ppt"}
_ODT_EXTENSIONS = {".odt", ".ods", ".odp"}
_RTF_EXTENSIONS = {".rtf"}
_TEXT_EXTENSIONS = {".txt", ".csv", ".json", ".xml", ".html", ".htm", ".md", ".log"}

SUPPORTED_EXTENSIONS: set[str] = (
    _PDF_EXTENSIONS
    | _IMAGE_EXTENSIONS
    | _DOCX_EXTENSIONS
    | _EXCEL_EXTENSIONS
    | _PPTX_EXTENSIONS
    | _ODT_EXTENSIONS
    | _RTF_EXTENSIONS
    | _TEXT_EXTENSIONS
)

# ── Dynamic import flags for format support ──
try:
    import docx as _docx
    _HAS_DOCX = True
except ImportError:
    _docx = None
    _HAS_DOCX = False

try:
    import openpyxl
    _HAS_OPENPYXL = True
except ImportError:
    openpyxl = None
    _HAS_OPENPYXL = False

try:
    import pillow_heif as _pillow_heif
    _pillow_heif.register_heif_opener()
    _HAS_HEIF = True
except Exception:
    _HAS_HEIF = False

try:
    from pptx import Presentation as _Presentation
    _HAS_PPTX = True
except ImportError:
    _Presentation = None
    _HAS_PPTX = False

try:
    from odf.opendocument import load as _odf_load
    from odf import text as _odf_text, teletype as _odf_teletype, table as _odf_table
    _HAS_ODF = True
except ImportError:
    _odf_load = None
    _HAS_ODF = False

try:
    from striprtf.striprtf import rtf_to_text as _rtf_to_text
    _HAS_STRIPRTF = True
except ImportError:
    _rtf_to_text = None
    _HAS_STRIPRTF = False

try:
    import pikepdf as _pikepdf
    _HAS_PIKEPDF = True
except ImportError:
    _pikepdf = None
    _HAS_PIKEPDF = False

import pdfplumber


def _extract_fields_and_tables(text: str) -> dict:
    """
    Dynamically extract key-value fields and tables from OCR text.
    Returns a dict: { 'fields': {key: value, ...}, 'tables': [table1, ...] }
    """
    fields = {}
    tables = []
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    kv_pattern = re.compile(r"^([A-Za-z0-9 .\-_/]+)\s*[:\-–]\s*(.+)$")
    for line in lines:
        m = kv_pattern.match(line)
        if m:
            key, value = m.group(1).strip(), m.group(2).strip()
            if key and value:
                fields[key] = value

    current_table = []
    for line in lines:
        if '|' in line:
            cols = [c.strip() for c in line.split('|')]
        else:
            cols = re.split(r"\s{2,}|\t", line)
        if len([c for c in cols if c]) >= 2:
            current_table.append(cols)
        else:
            if current_table:
                tables.append(current_table)
                current_table = []
    if current_table:
        tables.append(current_table)
    return {'fields': fields, 'tables': tables}


def _format_table(table: list) -> str:
    """Format a 2D grid/list of lists table as a markdown string."""
    lines = []
    for row in table:
        cells = [str(c).replace('\n', ' ').strip() for c in row]
        lines.append(" | ".join(cells))
    return "\n".join(lines)


def _detect_extractor_for_unknown(path: Path) -> str:
    try:
        with open(path, "rb") as fh:
            header = fh.read(8)
    except OSError:
        return "text"

    if header.startswith(b"%PDF"):
        return "pdf"
    if (
        header.startswith(b"\x89PNG")
        or header.startswith(b"\xff\xd8\xff")
        or header[:6] in (b"GIF87a", b"GIF89a")
        or header.startswith(b"BM")
        or header.startswith(b"II*\x00") or header.startswith(b"MM\x00*")
        or header[:4] == b"RIFF"
    ):
        return "image"
    if header[:2] == b"PK":
        return "docx"
    return "text"


def extract_with_azure_docintel(file_path: Path, document_id: str | None = None) -> list[dict] | None:
    endpoint = getattr(settings, "azure_document_intelligence_endpoint", None) or os.environ.get("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT")
    key = getattr(settings, "azure_document_intelligence_key", None) or os.environ.get("AZURE_DOCUMENT_INTELLIGENCE_KEY")
    if not endpoint or not key:
        logger.warning("Azure Document Intelligence credentials not configured. Skipping Azure OCR.")
        return None

    try:
        from azure.ai.formrecognizer import DocumentAnalysisClient
        from azure.core.credentials import AzureKeyCredential
    except ImportError:
        logger.warning("azure-ai-formrecognizer package is not installed. Skipping Azure OCR.")
        return None

    logger.info("Extracting layout from %s using Azure Document Intelligence...", file_path.name)
    try:
        client = DocumentAnalysisClient(endpoint=endpoint, credential=AzureKeyCredential(key))
        with open(file_path, "rb") as f:
            poller = client.begin_analyze_document("prebuilt-layout", document=f)
            result = poller.result()

        pages_data = []
        # Group tables by page
        tables_by_page = {}
        for table in getattr(result, "tables", []):
            for cell in table.cells:
                if cell.bounding_regions:
                    page_num = cell.bounding_regions[0].page_number
                    if page_num not in tables_by_page:
                        tables_by_page[page_num] = []
                    t_idx = result.tables.index(table)
                    if t_idx not in [t[0] for t in tables_by_page[page_num]]:
                        tables_by_page[page_num].append((t_idx, table))

        for page in result.pages:
            page_num = page.page_number
            page_lines = [line.content for line in getattr(page, "lines", [])]
            page_text = "\n".join(page_lines)

            # Reconstruct tables for this page
            page_tables = []
            formatted_table_texts = []
            for t_idx, table in tables_by_page.get(page_num, []):
                grid = [["" for _ in range(table.column_count)] for _ in range(table.row_count)]
                for cell in table.cells:
                    grid[cell.row_index][cell.column_index] = cell.content or ""
                page_tables.append(grid)
                formatted = _format_table(grid)
                if formatted:
                    formatted_table_texts.append(formatted)

            # Append formatted table text to the page text
            if formatted_table_texts:
                page_text += "\n\n" + "\n\n".join(formatted_table_texts)

            # Extract word-level tokens
            tokens = []
            for word in getattr(page, "words", []):
                points = word.polygon
                if points and len(points) >= 4:
                    try:
                        xs = [p.x if hasattr(p, "x") else p[0] for p in points]
                        ys = [p.y if hasattr(p, "y") else p[1] for p in points]
                        x0, y0, x1, y1 = min(xs), min(ys), max(xs), max(ys)
                    except Exception:
                        x0, y0, x1, y1 = 0.0, 0.0, 0.0, 0.0
                else:
                    x0, y0, x1, y1 = 0.0, 0.0, 0.0, 0.0

                tokens.append({
                    "text": word.content,
                    "x0": x0,
                    "y0": y0,
                    "x1": x1,
                    "y1": y1,
                    "page": page_num,
                    "confidence": word.confidence,
                    "source": "azure_docintel",
                })

            parsed = _extract_fields_and_tables(page_text)
            pages_data.append({
                'page': page_num,
                'text': page_text,
                'fields': parsed['fields'],
                'tables': page_tables,
                'confidence': sum(w.confidence for w in getattr(page, "words", [])) / len(getattr(page, "words", [])) if getattr(page, "words", None) else 1.0,
                'tokens': tokens,
            })

        logger.info("Successfully extracted %d pages from %s using Azure Document Intelligence", len(pages_data), file_path.name)
        
        # Dump output for pipeline analysis & verification
        try:
            dump_path = Path("/app/tmp")
            if not dump_path.exists():
                dump_path = Path("tmp")
            dump_path.mkdir(parents=True, exist_ok=True)
            if document_id:
                json_name = f"azure_ocr_response_{document_id}.json"
            else:
                json_name = f"azure_ocr_response_{file_path.name}.json"
            target_json = dump_path / json_name
            target_json.write_text(json.dumps(pages_data, indent=2, ensure_ascii=False), encoding="utf-8")
            logger.info("Dumped Azure Document Intelligence layout extraction to %s", target_json)
        except Exception as dump_err:
            logger.warning("Failed to dump Azure OCR output to tmp folder: %s", dump_err)

        return pages_data

    except Exception as e:
        logger.error("Azure Document Intelligence extraction failed: %s", e, exc_info=True)
        return None


def _maybe_decrypt_pdf(path: Path) -> Path:
    if not _HAS_PIKEPDF:
        return path

    try:
        with open(path, "rb") as fh:
            head = fh.read(1024)
        if b"/Encrypt" not in head and b"/Encrypt " not in head:
            return path
    except Exception:
        return path

    candidate_passwords: list[str] = [""]
    extra = os.environ.get("OCR_PDF_PASSWORDS", "") or getattr(settings, "pdf_passwords", "")
    for pw in (extra or "").split(","):
        pw = pw.strip()
        if pw and pw not in candidate_passwords:
            candidate_passwords.append(pw)

    last_err: Exception | None = None
    for pw in candidate_passwords:
        try:
            with _pikepdf.open(str(path), password=pw) as pdf:
                tmp = tempfile.NamedTemporaryFile(
                    delete=False, suffix=".pdf", prefix="ocr_decrypted_",
                )
                tmp.close()
                pdf.save(tmp.name)
                logger.info(
                    "Decrypted encrypted PDF (password=%s) -> %s",
                    "<empty>" if not pw else "<provided>", tmp.name,
                )
                return Path(tmp.name)
        except Exception as exc:
            last_err = exc
            continue

    logger.warning(
        "Could not decrypt PDF %s with any candidate password (set OCR_PDF_PASSWORDS); "
        "downstream parser may fail. Last error: %s",
        path, last_err,
    )
    return path


def _process_pdf_page_worker(pdf_path: str, page_idx: int) -> dict:
    """Extract text from a digital PDF page using pdfplumber."""
    with pdfplumber.open(pdf_path) as pdf:
        page = pdf.pages[page_idx]
        page_num = page_idx + 1
        
        parts: list[str] = []
        tables_found = []

        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())

        try:
            tables = page.extract_tables()
            for table in (tables or []):
                table_text = _format_table(table)
                if table_text:
                    parts.append(table_text)
                    tables_found.append(table)
        except Exception:
            logger.debug("Table extraction failed on page %d", page_num)

        digital_text = "\n\n".join(parts).strip()
        parsed = _extract_fields_and_tables(digital_text)
        
        tokens: list[dict[str, Any]] = []
        try:
            words = page.extract_words()
            for w in words:
                tokens.append({
                    "text": w.get("text", "").strip(),
                    "x0": float(w.get("x0", 0)),
                    "y0": float(w.get("top", w.get("y0", 0))),
                    "x1": float(w.get("x1", 0)),
                    "y1": float(w.get("bottom", w.get("y1", 0))),
                    "page": page_num,
                    "confidence": None,
                    "source": "pdfplumber",
                })
        except Exception:
            pass

        if tables_found:
            parsed['tables'] = tables_found + parsed['tables']
        
        return {
            'page': page_num,
            'text': digital_text,
            'fields': parsed['fields'],
            'tables': parsed['tables'],
            'confidence': 1.0,
            'tokens': tokens,
        }


def _extract_from_pdf(path: Path) -> list[dict]:
    working_path = _maybe_decrypt_pdf(path)
    with pdfplumber.open(working_path) as pdf:
        num_pages = len(pdf.pages)
    
    max_workers = min(4, max(2, num_pages // 2))
    page_indices = list(range(num_pages))
    
    results: list[dict] = []
    try:
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            page_results = executor.map(
                _process_pdf_page_worker,
                [str(working_path)] * num_pages,
                page_indices
            )
            results = list(page_results)
    except Exception as e:
        logger.error("Parallel PDF page processing failed, falling back to sequential: %s", e)
        for page_idx in page_indices:
            try:
                result = _process_pdf_page_worker(str(working_path), page_idx)
                results.append(result)
            except Exception as page_err:
                logger.error("Error processing page %d: %s", page_idx + 1, page_err)
    
    return results


def _extract_from_image(path: Path) -> list[dict]:
    azure_res = extract_with_azure_docintel(path)
    if azure_res is not None:
        return azure_res
    raise ValueError("Azure Document Intelligence failed to extract text from image.")


def _extract_from_docx(path: Path) -> list[dict]:
    if not _HAS_DOCX:
        raise ValueError("python-docx not installed -- cannot process .docx files")

    doc = _docx.Document(str(path))
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))

    combined = "\n\n".join(parts)
    parsed = _extract_fields_and_tables(combined)
    return [{
        'page': 1,
        'text': combined,
        'fields': parsed['fields'],
        'tables': parsed['tables'],
        'confidence': 1.0,
        'tokens': []
    }]


def _extract_from_excel(path: Path) -> list[dict]:
    if not _HAS_OPENPYXL:
        raise ValueError("openpyxl not installed -- cannot process .xlsx files")

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    results: list[dict] = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        lines: list[str] = []
        lines.append(f"[Sheet: {sheet_name}]")

        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                lines.append(" | ".join(cells))

        text = "\n".join(lines)
        parsed = _extract_fields_and_tables(text)
        results.append({
            'page': sheet_idx,
            'text': text,
            'fields': parsed['fields'],
            'tables': parsed['tables'],
            'confidence': 1.0,
            'tokens': []
        })

    wb.close()
    return results if results else [{
        'page': 1,
        'text': "",
        'fields': {},
        'tables': [],
        'confidence': None,
        'tokens': []
    }]


def _extract_from_pptx(path: Path) -> list[dict]:
    if not _HAS_PPTX:
        raise ValueError("python-pptx not installed -- cannot process .pptx files")

    prs = _Presentation(str(path))
    results: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                parts.append(f"[Speaker notes]\n{notes.strip()}")
        except Exception:
            pass
        text = "\n".join(parts).strip()
        parsed = _extract_fields_and_tables(text)
        results.append({
            'page': slide_idx,
            'text': text,
            'fields': parsed['fields'],
            'tables': parsed['tables'],
            'confidence': 1.0,
            'tokens': []
        })
    return results if results else [{
        'page': 1,
        'text': "",
        'fields': {},
        'tables': [],
        'confidence': None,
        'tokens': []
    }]


def _extract_from_odt(path: Path) -> list[dict]:
    if not _HAS_ODF:
        raise ValueError("odfpy not installed -- cannot process OpenDocument files")

    doc = _odf_load(str(path))
    parts: list[str] = []

    for p in doc.getElementsByType(_odf_text.P):
        line = _odf_teletype.extractText(p).strip()
        if line:
            parts.append(line)

    for tbl in doc.getElementsByType(_odf_table.Table):
        for row in tbl.getElementsByType(_odf_table.TableRow):
            cells = []
            for cell in row.getElementsByType(_odf_table.TableCell):
                cells.append(_odf_teletype.extractText(cell).strip())
            if any(cells):
                parts.append(" | ".join(cells))

    text = "\n".join(parts).strip()
    parsed = _extract_fields_and_tables(text)
    return [{
        'page': 1,
        'text': text,
        'fields': parsed['fields'],
        'tables': parsed['tables'],
        'confidence': 1.0,
        'tokens': []
    }]


def _extract_from_rtf(path: Path) -> list[dict]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    if _HAS_STRIPRTF and _rtf_to_text is not None:
        try:
            text = _rtf_to_text(raw, errors="ignore").strip()
        except Exception:
            text = raw
    else:
        text = re.sub(r"\\\w+\b", "", raw)
        text = text.replace("{", "").replace("}", "")
        text = re.sub(r"\s+", " ", text).strip()

    parsed = _extract_fields_and_tables(text)
    return [{
        'page': 1,
        'text': text,
        'fields': parsed['fields'],
        'tables': parsed['tables'],
        'confidence': 1.0,
        'tokens': []
    }]


def _extract_from_text(path: Path) -> list[dict]:
    suffix = path.suffix.lower()

    for encoding in ("utf-8", "latin-1"):
        try:
            raw = path.read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, ValueError):
            continue
    else:
        raw = path.read_bytes().decode("utf-8", errors="replace")

    if suffix == ".csv":
        return _extract_from_csv_text(raw)
    if suffix == ".json":
        return _extract_from_json_text(raw)

    parsed = _extract_fields_and_tables(raw.strip())
    return [{
        'page': 1,
        'text': raw.strip(),
        'fields': parsed['fields'],
        'tables': parsed['tables'],
        'confidence': 1.0,
        'tokens': []
    }]


def _extract_from_csv_text(raw: str) -> list[dict]:
    reader = csv.reader(io.StringIO(raw))
    lines: list[str] = []
    for row in reader:
        cells = [c.strip() for c in row]
        if any(cells):
            lines.append(" | ".join(cells))
    text = "\n".join(lines)
    parsed = _extract_fields_and_tables(text)
    return [{
        'page': 1,
        'text': text,
        'fields': parsed['fields'],
        'tables': parsed['tables'],
        'confidence': 1.0,
        'tokens': []
    }]


def _extract_from_json_text(raw: str) -> list[dict]:
    try:
        data = json.loads(raw)
        text = json.dumps(data, indent=2)
    except Exception:
        text = raw.strip()
    parsed = _extract_fields_and_tables(text)
    return [{
        'page': 1,
        'text': text,
        'fields': parsed['fields'],
        'tables': parsed['tables'],
        'confidence': 1.0,
        'tokens': []
    }]


def extract_text(file_path: str | Path, document_id: str | None = None) -> list[dict]:
    """Run extraction on any supported file using Azure OCR or digital extraction."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    # Azure Document Intelligence integration (for scanned files / images)
    if (suffix in _PDF_EXTENSIONS or suffix in _IMAGE_EXTENSIONS):
        is_digital_pdf = False
        if suffix in _PDF_EXTENSIONS:
            try:
                import pdfplumber
                with pdfplumber.open(path) as pdf:
                    total_len = 0
                    for page in pdf.pages[:3]:
                        text = page.extract_text() or ""
                        total_len += len(re.sub(r"\s+", "", text))
                    if total_len >= 20:
                        is_digital_pdf = True
                        logger.info("PDF %s has extractable digital text (len=%d >= 20). Bypassing Azure OCR to optimize cost.", path.name, total_len)
            except Exception:
                pass

        if not is_digital_pdf:
            azure_res = extract_with_azure_docintel(path, document_id=document_id)
            if azure_res is not None:
                return azure_res
            else:
                raise ValueError("Azure Document Intelligence failed to extract text from scanned file.")

    if suffix in _PDF_EXTENSIONS:
        return _extract_from_pdf(path)
    elif suffix in _IMAGE_EXTENSIONS:
        return _extract_from_image(path)
    elif suffix in _DOCX_EXTENSIONS:
        return _extract_from_docx(path)
    elif suffix in _EXCEL_EXTENSIONS:
        return _extract_from_excel(path)
    elif suffix in _PPTX_EXTENSIONS:
        return _extract_from_pptx(path)
    elif suffix in _ODT_EXTENSIONS:
        return _extract_from_odt(path)
    elif suffix in _RTF_EXTENSIONS:
        return _extract_from_rtf(path)
    elif suffix in _TEXT_EXTENSIONS:
        return _extract_from_text(path)
    else:
        kind = _detect_extractor_for_unknown(path)
        attempts = {
            "pdf": _extract_from_pdf,
            "image": _extract_from_image,
            "docx": _extract_from_docx,
            "excel": _extract_from_excel,
            "text": _extract_from_text,
        }
        order = [kind] + [k for k in ("image", "pdf", "docx", "excel", "text") if k != kind]
        for k in order:
            try:
                return attempts[k](path)
            except Exception:
                continue
        raise ValueError(f"Failed to extract text from file: {path.name}")


def prewarm_ocr_engines() -> None:
    """Pre-warming is handled serverless by Azure Document Intelligence; no local prewarming needed."""
    pass


extract_text_structured = extract_text
